from __future__ import annotations

import base64
import io
import json
from typing import Any


EXECUTIVE_EXPORT_TERMS = (
    "Executive Brief",
    "História Executiva",
    "Scorecard Executivo",
    "Top 3 Drivers",
    "Decisões Prioritárias",
    "Plano de Ação",
)


def validate_report_quality(
    report: dict[str, Any],
    *,
    exports: dict[str, dict[str, Any]] | None = None,
) -> dict[str, Any]:
    checks = _build_checks(report, exports or {})
    score = _decision_value_score(report, checks)
    failures = _failures(report, checks, score)
    warnings = _warnings(checks, score)
    status = "FAIL" if failures else "WARNING" if warnings else "PASS"
    return {
        "status": status,
        "decision_value_score": score,
        "classification": _classification(score),
        "checks": checks,
        "failures": failures,
        "warnings": warnings,
        "stage_comparison": build_stage_comparison(report),
        "export_validation": validate_exports(exports or {}),
        "evidence_summary": _evidence_summary(report),
    }


def build_stage_comparison(report: dict[str, Any]) -> list[dict[str, Any]]:
    data = report.get("data", {})
    analytical = report.get("analytical", {})
    narrative = report.get("executive_narrative", {})
    discovery = report.get("discovery", {})
    story = report.get("executive_story", {})
    base_metrics = report.get("metrics", {})

    stages = [
        ("base", {"data": data, "metrics": base_metrics}),
        ("analytical_enrichment", {"analytical": analytical}),
        ("executive_narrative", {"executive_narrative": narrative}),
        ("discovery", {"discovery": discovery}),
        ("executive_story", {"executive_story": story}),
    ]
    result = []
    for name, payload in stages:
        merged = _merge_until(report, name)
        result.append(
            {
                "stage": name,
                "metric_count": _count_metrics(payload),
                "evidence_count": _count_evidence(payload),
                "decision_count": _count_decisions(payload),
                "actionable_recommendation_count": _count_recommendations(payload),
                "decision_value_score": _decision_value_score(merged, _build_checks(merged, {})),
                "report_quality_score": _score_value(merged.get("report_quality_score")),
                "report_intelligence_score": _score_value(merged.get("report_intelligence_score")),
                "executive_story_quality_score": _score_value(merged.get("executive_story_quality_score")),
            }
        )
    return result


def validate_exports(exports: dict[str, dict[str, Any]]) -> dict[str, Any]:
    expected = ("json", "markdown", "pdf", "pptx")
    checks: dict[str, Any] = {}
    for fmt in expected:
        payload = exports.get(fmt)
        if not payload:
            checks[fmt] = {
                "present": False,
                "carries_executive_blocks": False,
                "evidence": ["export ausente"],
            }
            continue
        content = _decode_export(payload)
        carries = _export_carries_blocks(fmt, payload, content)
        checks[fmt] = {
            "present": True,
            "carries_executive_blocks": carries,
            "content_type": payload.get("content_type"),
            "size_bytes": payload.get("size_bytes", 0),
            "evidence": _export_evidence(fmt, payload, content, carries),
        }
    return checks


def _build_checks(report: dict[str, Any], exports: dict[str, dict[str, Any]]) -> dict[str, Any]:
    story = report.get("executive_story") or {}
    analytical = report.get("analytical") or {}
    narrative = report.get("executive_narrative") or {}
    discovery = report.get("discovery") or {}
    metrics_pack = analytical.get("metrics_pack") or {}
    risks = metrics_pack.get("risks") or {}
    recommendations = analytical.get("recommendations") or []
    export_validation = validate_exports(exports)

    checks = {
        "executive_story_present": bool(story.get("generated") and story.get("period_story")),
        "top_3_drivers_present": 1 <= len(story.get("key_drivers", [])) <= 3,
        "decisions_with_evidence": _all_have_evidence(story.get("decision_ready_summary", []), "evidence"),
        "actionable_recommendations": any(item.get("action") and item.get("evidence") for item in recommendations),
        "risks_prioritized": bool(risks.get("high_risk_cards") or narrative.get("insights")),
        "analytical_metrics_present": bool(metrics_pack),
        "discovery_present": bool(discovery.get("report_intelligence_score")),
        "description_considered": _description_considered(report),
        "no_empty_critical_sections": _no_empty_critical_sections(story),
        "exports_carry_narrative": all(
            item.get("carries_executive_blocks")
            for key, item in export_validation.items()
            if key in {"json", "markdown", "pdf", "pptx"}
        ) if export_validation else False,
        "forecast_present": bool(discovery.get("what_happens_next")),
        "evidence_traceable": bool(story.get("evidence_map")),
    }
    return checks


def _decision_value_score(report: dict[str, Any], checks: dict[str, Any]) -> int:
    story_quality = _score_value(report.get("executive_story_quality_score"))
    readability = _score_value(report.get("executive_readability_score"))
    intelligence = _score_value(report.get("report_intelligence_score"))
    report_quality = _score_value(report.get("report_quality_score"))
    weighted = {
        "clareza_executiva": 15 if checks.get("executive_story_present") else 0,
        "capacidade_decisao": 15 if checks.get("decisions_with_evidence") else 0,
        "evidencia": 15 if checks.get("evidence_traceable") else 0,
        "priorizacao": 10 if checks.get("top_3_drivers_present") and checks.get("risks_prioritized") else 0,
        "profundidade_analitica": min(15, int((report_quality + intelligence) / 200 * 15)),
        "recomendacoes_acionaveis": 10 if checks.get("actionable_recommendations") else 0,
        "previsao_cenario": 10 if checks.get("forecast_present") else 0,
        "qualidade_exports": 10 if checks.get("exports_carry_narrative") else 0,
    }
    base = sum(weighted.values())
    quality_bonus = int((story_quality + readability) / 200 * 10)
    return min(100, base + quality_bonus)


def _classification(score: int) -> str:
    if score >= 80:
        return "executivo"
    if score >= 60:
        return "bom"
    if score >= 40:
        return "aceitavel"
    return "fraco"


def _failures(report: dict[str, Any], checks: dict[str, Any], score: int) -> list[str]:
    failures = []
    if score < 70:
        failures.append(f"DecisionValueScore abaixo do gate: {score} < 70")
    if not checks.get("executive_story_present"):
        failures.append("executive_story ausente ou nao gerado")
    if not checks.get("decisions_with_evidence"):
        failures.append("decisoes sem evidencia")
    if not checks.get("actionable_recommendations"):
        failures.append("relatorio sem recomendacoes acionaveis")
    if not checks.get("exports_carry_narrative"):
        failures.append("exports nao carregam blocos executivos")
    return failures


def _warnings(checks: dict[str, Any], score: int) -> list[str]:
    warnings = []
    if score < 80 and score >= 70:
        warnings.append("DecisionValueScore passou no gate, mas ainda nao atingiu nivel executivo")
    for key, passed in checks.items():
        if not passed and key not in {
            "executive_story_present",
            "decisions_with_evidence",
            "actionable_recommendations",
            "exports_carry_narrative",
        }:
            warnings.append(f"criterio nao atendido: {key}")
    return warnings


def _description_considered(report: dict[str, Any]) -> bool:
    quality = ((report.get("analytical") or {}).get("metrics_pack") or {}).get("quality") or {}
    cards = ((report.get("analytical") or {}).get("activity_classification") or {}).get("cards") or []
    return "incomplete_description_count" in quality or any("description_quality_score" in card for card in cards)


def _no_empty_critical_sections(story: dict[str, Any]) -> bool:
    structure = story.get("story_structure") or {}
    critical = (
        "contexto_periodo",
        "fatores_resultado",
        "decisoes_prioritarias",
        "plano_acao",
    )
    return all(
        structure.get(key, {}).get("summary") and structure.get(key, {}).get("evidence")
        for key in critical
    )


def _all_have_evidence(items: list[dict[str, Any]], field: str) -> bool:
    return bool(items) and all(item.get(field) for item in items)


def _evidence_summary(report: dict[str, Any]) -> dict[str, int]:
    return {
        "evidence_map": len((report.get("executive_story") or {}).get("evidence_map", [])),
        "driver_evidence": _count_evidence((report.get("executive_story") or {}).get("key_drivers", [])),
        "decision_evidence": _count_evidence((report.get("executive_story") or {}).get("decision_ready_summary", [])),
        "recommendation_evidence": _count_evidence(((report.get("analytical") or {}).get("recommendations") or [])),
    }


def _merge_until(report: dict[str, Any], stage: str) -> dict[str, Any]:
    merged = {
        "data": report.get("data", {}),
        "metrics": report.get("metrics", {}),
    }
    if stage in {"analytical_enrichment", "executive_narrative", "discovery", "executive_story"}:
        merged["analytical"] = report.get("analytical", {})
        merged["report_quality_score"] = report.get("report_quality_score")
    if stage in {"executive_narrative", "discovery", "executive_story"}:
        merged["executive_narrative"] = report.get("executive_narrative", {})
        merged["executive_readability_score"] = report.get("executive_readability_score")
    if stage in {"discovery", "executive_story"}:
        merged["discovery"] = report.get("discovery", {})
        merged["report_intelligence_score"] = report.get("report_intelligence_score")
    if stage == "executive_story":
        merged["executive_story"] = report.get("executive_story", {})
        merged["executive_story_quality_score"] = report.get("executive_story_quality_score")
    return merged


def _count_metrics(value: Any) -> int:
    if isinstance(value, dict):
        return sum(1 for v in value.values() if isinstance(v, (int, float, str, list, dict)) and v not in (None, "", [], {})) + sum(_count_metrics(v) for v in value.values())
    if isinstance(value, list):
        return sum(_count_metrics(item) for item in value)
    return 0


def _count_evidence(value: Any) -> int:
    if isinstance(value, dict):
        total = 0
        for key, item in value.items():
            if key in {"evidence", "evidence_map", "reasons", "basis"}:
                total += len(item) if isinstance(item, list) else int(bool(item))
            total += _count_evidence(item)
        return total
    if isinstance(value, list):
        return sum(_count_evidence(item) for item in value)
    return 0


def _count_decisions(value: Any) -> int:
    if isinstance(value, dict):
        total = 0
        for key, item in value.items():
            if key in {"decision_ready_summary", "priority_decisions", "management_decisions"} and isinstance(item, list):
                total += len(item)
            total += _count_decisions(item)
        return total
    if isinstance(value, list):
        return sum(_count_decisions(item) for item in value)
    return 0


def _count_recommendations(value: Any) -> int:
    if isinstance(value, dict):
        total = 0
        for key, item in value.items():
            if key in {"recommendations", "opportunities", "action_plan"} and isinstance(item, list):
                total += len([entry for entry in item if entry.get("action") or entry.get("recommended_action") or entry.get("title")])
            total += _count_recommendations(item)
        return total
    if isinstance(value, list):
        return sum(_count_recommendations(item) for item in value)
    return 0


def _score_value(value: Any) -> int:
    if isinstance(value, dict):
        return int(value.get("score") or value.get("report_quality_score") or 0)
    if isinstance(value, (int, float)):
        return int(value)
    return 0


def _decode_export(payload: dict[str, Any]) -> str:
    content = payload.get("content_base64")
    if not content:
        return ""
    try:
        data = base64.b64decode(content)
        return data.decode("utf-8-sig", errors="ignore")
    except Exception:
        return ""


def _export_carries_blocks(fmt: str, payload: dict[str, Any], content: str) -> bool:
    if fmt == "json":
        return bool(payload.get("executive_story") or payload.get("discovery") or payload.get("executive_narrative"))
    if fmt == "markdown":
        return all(term in content for term in EXECUTIVE_EXPORT_TERMS)
    if fmt == "pptx":
        try:
            data = json.loads(content)
            slide_titles = [slide.get("title") for slide in data.get("slides", [])]
            return all(title in slide_titles for title in EXECUTIVE_EXPORT_TERMS)
        except json.JSONDecodeError:
            slide_titles = _pptx_slide_titles(payload)
        return all(title in slide_titles for title in EXECUTIVE_EXPORT_TERMS)
    if fmt == "pdf":
        return payload.get("content_type") == "application/pdf" and int(payload.get("size_bytes") or 0) > 0
    return False


def _export_evidence(fmt: str, payload: dict[str, Any], content: str, carries: bool) -> list[str]:
    if fmt == "pdf":
        return [f"content_type={payload.get('content_type')}", f"size_bytes={payload.get('size_bytes')}"]
    if fmt == "json":
        found = [
            key
            for key in ("executive_story", "discovery", "executive_narrative")
            if payload.get(key)
        ]
        return [f"{key} presente no payload" for key in found] or ["blocos executivos ausentes no payload"]
    if fmt == "pptx" and "PK" not in content[:4]:
        found = [term for term in EXECUTIVE_EXPORT_TERMS if term in content]
        return found or ["blocos executivos nao encontrados"]
    if fmt == "pptx":
        found = [title for title in _pptx_slide_titles(payload) if title in EXECUTIVE_EXPORT_TERMS]
        return found or ["slides executivos nao encontrados"]
    found = [term for term in EXECUTIVE_EXPORT_TERMS if term in content]
    return found or ["blocos executivos nao encontrados"]


def _pptx_slide_titles(payload: dict[str, Any]) -> list[str]:
    try:
        from pptx import Presentation

        raw = base64.b64decode(payload.get("content_base64") or "")
        presentation = Presentation(io.BytesIO(raw))
        titles = []
        for slide in presentation.slides:
            if slide.shapes.title and slide.shapes.title.text:
                titles.append(slide.shapes.title.text.strip())
                continue
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                    titles.append(shape.text_frame.text.strip())
                    break
        return titles
    except Exception:
        return []
